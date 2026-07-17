"""
Build a professional PPTX report from results_summary.json
Output: EMA_Stack_Backtest_Results.pptx
"""
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette ----
DARK   = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
ACCENT = RGBColor(0x00, 0xC8, 0x96)   # teal/green
ACCENT2= RGBColor(0x3B, 0x82, 0xF6)   # blue
RED    = RGBColor(0xEF, 0x44, 0x44)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x94, 0xA3, 0xB8)
LIGHT  = RGBColor(0xE2, 0xE8, 0xF0)
CARD   = RGBColor(0x1E, 0x29, 0x3B)

SW, SH = Inches(13.333), Inches(7.5)

res = json.load(open("results_summary.json"))
prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def bg(slide, color=DARK):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h, color):
    sp = slide.shapes.add_shape(1, l, t, w, h)  # rectangle
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def text(slide, l, t, w, h, s, size=18, color=WHITE, bold=False,
         align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4)
    lines = s.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def bar(slide, l, t, maxw, pct, color, h=Inches(0.28)):
    """horizontal win-rate bar; pct 0..100"""
    box(slide, l, t, maxw, h, CARD)
    w = int(maxw * min(max(pct, 0), 100) / 100)
    if w > 0:
        box(slide, l, t, Emu(w), h, color)


# ============================================================ SLIDE 1: TITLE
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0, Inches(3.05), SW, Inches(0.06), ACCENT)
text(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.2),
     "EMA 8/24/72 Stack Strategy", size=48, bold=True, color=WHITE)
text(s, Inches(0.8), Inches(2.35), Inches(11.7), Inches(0.7),
     "Backtest Results — NIFTY 50, Bank Nifty & ATM Options", size=24, color=ACCENT)
text(s, Inches(0.8), Inches(3.4), Inches(11.7), Inches(1.0),
     "Real 5-minute & daily market data from Dhan API\n"
     "Includes rare FREE historical OPTIONS data (rolling ATM Call/Put)",
     size=16, color=GREY)
text(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5),
     "Data source: DhanHQ  |  Strategy: EMA stack crossover + volume confirmation",
     size=12, color=GREY)

# ============================================================ SLIDE 2: WHAT IS IT
s = prs.slides.add_slide(BLANK)
bg(s)
text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
     "The Strategy", size=34, bold=True, color=ACCENT)
cards = [
    ("EMA Stack", "3 EMAs: 8 / 24 / 72.\nBULLISH when 8>24>72.\nBEARISH when 8<24<72.", ACCENT2),
    ("Fresh Signal", "Trade fires only on the\ncandle the stack first\nforms (not every bar).", ACCENT),
    ("Volume Filter", "Signal must have volume\n>= 1.3x its 20-bar\naverage to confirm.", RGBColor(0xF5,0x9E,0x0B)),
    ("Horizons", "Forward return checked at\n+3/+6/+12/+24 bars\n(or trading days).", ACCENT2),
]
x = Inches(0.7)
for title, body, col in cards:
    box(s, x, Inches(1.7), Inches(2.85), Inches(3.2), CARD)
    box(s, x, Inches(1.7), Inches(2.85), Inches(0.12), col)
    text(s, x+Inches(0.2), Inches(2.0), Inches(2.5), Inches(0.6), title, size=20, bold=True, color=col)
    text(s, x+Inches(0.2), Inches(2.75), Inches(2.5), Inches(2.0), body, size=15, color=LIGHT)
    x += Inches(3.05)
text(s, Inches(0.7), Inches(5.5), Inches(12), Inches(1.5),
     "Tested on: Nifty 50 index, Bank Nifty index, all 48 liquid Nifty-50 stocks (3.5 yrs daily),\n"
     "and rolling ATM Call/Put options for Nifty & Bank Nifty (6 months of 5-min data).",
     size=16, color=GREY)

# ============================================================ helper: ranking slide
def ranking_slide(title, subtitle, rows, note=""):
    s = prs.slides.add_slide(BLANK)
    bg(s)
    text(s, Inches(0.7), Inches(0.45), Inches(12), Inches(0.7),
         title, size=32, bold=True, color=ACCENT)
    text(s, Inches(0.7), Inches(1.15), Inches(12), Inches(0.5),
         subtitle, size=15, color=GREY)
    # header
    y = Inches(1.85)
    text(s, Inches(0.7),  y, Inches(2.6), Inches(0.35), "SYMBOL", size=13, bold=True, color=GREY)
    text(s, Inches(3.2),  y, Inches(1.2), Inches(0.35), "SIDE",   size=13, bold=True, color=GREY)
    text(s, Inches(4.3),  y, Inches(1.2), Inches(0.35), "SIGNALS",size=13, bold=True, color=GREY)
    text(s, Inches(5.6),  y, Inches(4.2), Inches(0.35), "WIN RATE",size=13, bold=True, color=GREY)
    text(s, Inches(10.0), y, Inches(2.5), Inches(0.35), "AVG RETURN",size=13, bold=True, color=GREY)
    y = Inches(2.3)
    rh = Inches(0.46)
    for sym, side, n, wr, avg in rows:
        col = ACCENT if wr >= 60 else (RGBColor(0xF5,0x9E,0x0B) if wr >= 50 else RED)
        text(s, Inches(0.7), y, Inches(2.6), rh, sym, size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(3.2), y, Inches(1.2), rh, side.upper(), size=13, color=LIGHT, anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(4.3), y, Inches(1.2), rh, str(n), size=13, color=LIGHT, anchor=MSO_ANCHOR.MIDDLE)
        bar(s, Inches(5.6), y+Inches(0.08), Inches(3.6), wr, col)
        text(s, Inches(9.25), y, Inches(0.9), rh, f"{wr:.0f}%", size=14, bold=True, color=col, anchor=MSO_ANCHOR.MIDDLE)
        ac = ACCENT if avg > 0 else RED
        text(s, Inches(10.0), y, Inches(2.5), rh, f"{avg:+.2f}%", size=14, bold=True, color=ac, anchor=MSO_ANCHOR.MIDDLE)
        y += rh
    if note:
        text(s, Inches(0.7), Inches(6.9), Inches(12), Inches(0.5), note, size=12, color=GREY)
    return s


# ---- build daily-spot ranking (+6d, min 5 signals) ----
rank = []
for sym, r in res["daily_spot"].items():
    for side in ("bull", "bear"):
        st = r[side].get("6")
        if st and st["signals"] >= 5:
            rank.append((sym, side, st["signals"], st["win_rate"], st["avg_ret"]))
rank.sort(key=lambda x: (x[3], x[4]), reverse=True)

ranking_slide(
    "Top Performers — Stocks & Indices",
    "Daily candles, 3.5 years of data. Horizon = +6 trading days. Minimum 5 signals.",
    rank[:10],
    "Green = win rate >= 60%. These are the setups with a real statistical edge.")

# ============================================================ SLIDE: INDEX SPOTLIGHT
s = prs.slides.add_slide(BLANK)
bg(s)
text(s, Inches(0.7), Inches(0.45), Inches(12), Inches(0.7),
     "Index Spotlight — Most Reliable Signals", size=30, bold=True, color=ACCENT)
text(s, Inches(0.7), Inches(1.15), Inches(12), Inches(0.5),
     "Indices have the largest signal counts, so their win rates are the most trustworthy.",
     size=15, color=GREY)

def big_stat(slide, l, name, side, st, col):
    box(slide, l, Inches(2.0), Inches(5.6), Inches(3.6), CARD)
    box(slide, l, Inches(2.0), Inches(5.6), Inches(0.14), col)
    text(slide, l+Inches(0.35), Inches(2.3), Inches(5), Inches(0.6),
         f"{name}  —  {side}", size=22, bold=True, color=WHITE)
    text(slide, l+Inches(0.35), Inches(3.05), Inches(3), Inches(1.4),
         f"{st['win_rate']:.0f}%", size=60, bold=True, color=col)
    text(slide, l+Inches(0.35), Inches(4.4), Inches(5), Inches(0.5),
         "Win rate  (+6 trading days)", size=14, color=GREY)
    text(slide, l+Inches(0.35), Inches(4.9), Inches(5), Inches(0.5),
         f"Signals: {st['signals']}     Avg return: {st['avg_ret']:+.2f}%",
         size=15, color=LIGHT)

bn = res["daily_spot"]["BANKNIFTY"]["bull"]["6"]
nf = res["daily_spot"]["NIFTY"]["bull"]["6"]
big_stat(s, Inches(0.7), "BANK NIFTY", "BULL", bn, ACCENT)
big_stat(s, Inches(7.0), "NIFTY 50", "BULL", nf, ACCENT)
text(s, Inches(0.7), Inches(6.1), Inches(12), Inches(1.0),
     "Verdict: On the BUY side, the EMA stack shows a genuine edge on indices\n"
     "(Bank Nifty 75% over 16 signals, Nifty 71% over 14 signals).",
     size=16, color=WHITE, bold=True)

# ============================================================ SLIDE: OPTIONS
s = prs.slides.add_slide(BLANK)
bg(s)
text(s, Inches(0.7), Inches(0.45), Inches(12), Inches(0.7),
     "Options Backtest (ATM Call / Put)", size=30, bold=True, color=ACCENT2)
text(s, Inches(0.7), Inches(1.15), Inches(12.0), Inches(0.6),
     "Rolling ATM options, 6 months of 5-min data. FREE data most sites don't give away.",
     size=15, color=GREY)
# header
y = Inches(1.95)
text(s, Inches(0.7), y, Inches(3.3), Inches(0.35), "OPTION LEG", size=13, bold=True, color=GREY)
text(s, Inches(4.0), y, Inches(1.6), Inches(0.35), "BEST WIN%", size=13, bold=True, color=GREY)
text(s, Inches(5.7), y, Inches(2.0), Inches(0.35), "SIGNALS", size=13, bold=True, color=GREY)
text(s, Inches(7.7), y, Inches(3.0), Inches(0.35), "AVG RET (+120min)", size=13, bold=True, color=GREY)
y = Inches(2.45)
rh = Inches(0.55)
for leg in ("BANKNIFTY_PUT", "BANKNIFTY_CALL", "NIFTY_PUT", "NIFTY_CALL"):
    o = res["options"].get(leg)
    if not o:
        continue
    # best of bull/bear at +24
    b, be = o["bull"]["24"], o["bear"]["24"]
    best = max([x for x in (b, be) if x], key=lambda z: z["win_rate"])
    side = "BEAR" if best is be else "BULL"
    col = ACCENT if best["win_rate"] >= 55 else (RGBColor(0xF5,0x9E,0x0B) if best["win_rate"] >= 45 else RED)
    text(s, Inches(0.7), y, Inches(3.3), rh, f"{leg.replace('_',' ')} ({side})", size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(4.0), y, Inches(1.6), rh, f"{best['win_rate']:.0f}%", size=17, bold=True, color=col, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(5.7), y, Inches(2.0), rh, str(best["signals"]), size=14, color=LIGHT, anchor=MSO_ANCHOR.MIDDLE)
    ac = ACCENT if best["avg_ret"] > 0 else RED
    text(s, Inches(7.7), y, Inches(3.0), rh, f"{best['avg_ret']:+.2f}%", size=15, bold=True, color=ac, anchor=MSO_ANCHOR.MIDDLE)
    y += rh
box(s, Inches(0.7), Inches(5.3), Inches(11.9), Inches(1.5), CARD)
box(s, Inches(0.7), Inches(5.3), Inches(0.14), Inches(1.5), RED)
text(s, Inches(1.0), Inches(5.5), Inches(11.4), Inches(1.2),
     "WARNING: On options, this strategy is WEAK. Buying options (BULL) loses on almost every leg\n"
     "(win rate 12-33%, negative avg return) due to theta decay + whipsaw. Only Bank Nifty PUT / 2-hr\n"
     "holds show ~59% — and even that swings -15% mid-trade. Use signals on SPOT/FUTURES, not options.",
     size=14, color=WHITE)

# ============================================================ SLIDE: VERDICT
s = prs.slides.add_slide(BLANK)
bg(s)
text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
     "Verdict & Honest Takeaways", size=34, bold=True, color=ACCENT)
points = [
    ("EDGE on indices (BUY side)", "Bank Nifty 75% & Nifty 71% win rate over 6 days. Largest, most reliable samples.", ACCENT),
    ("Some strong stocks", "TATASTEEL, BAJAJFINSV, ICICIBANK, BPCL show 70-100% — but smaller samples (5-7 signals).", ACCENT2),
    ("SHORT side is weak", "Bearish signals mostly break even or lose. Do not trade the short side blindly.", RED),
    ("Options = avoid", "Theta decay kills option buyers. Take the signal on spot/futures; use options only to execute.", RED),
    ("Data limitation", "Intraday 5-min is capped at ~5 days by Dhan; daily test covers 3.5 yrs. More data = more confidence.", GREY),
]
y = Inches(1.6)
for title, body, col in points:
    box(s, Inches(0.7), y, Inches(0.14), Inches(0.9), col)
    text(s, Inches(1.0), y, Inches(11.4), Inches(0.45), title, size=19, bold=True, color=col)
    text(s, Inches(1.0), y+Inches(0.42), Inches(11.4), Inches(0.5), body, size=14, color=LIGHT)
    y += Inches(1.05)

# ============================================================ SLIDE: CTA / DATA
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0, Inches(3.05), SW, Inches(0.06), ACCENT)
text(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.0),
     "Free Data & Full Backtest Code", size=40, bold=True, color=WHITE)
text(s, Inches(0.8), Inches(2.5), Inches(11.7), Inches(0.6),
     "Everything used in this study is shared openly.", size=20, color=ACCENT)
text(s, Inches(0.8), Inches(3.4), Inches(11.7), Inches(2.5),
     "• 5-min & daily OHLCV for Nifty 50, Bank Nifty and 48 stocks\n"
     "• Rare FREE historical OPTIONS data — rolling ATM Call/Put (6 months)\n"
     "• Python backtest scripts (EMA stack strategy, fully documented)\n"
     "• Excel workbooks with every candle + a ranking sheet",
     size=18, color=LIGHT)
text(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6),
     "Search terms: nifty 5 min data free · bank nifty option historical data · "
     "EMA crossover backtest India · dhan api backtest",
     size=12, color=GREY)

prs.save("EMA_Stack_Backtest_Results.pptx")
print("Saved EMA_Stack_Backtest_Results.pptx with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
